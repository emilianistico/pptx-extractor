
from fastapi import FastAPI
from pydantic import BaseModel
import base64
import io
from pptx import Presentation
from fastapi.responses import JSONResponse

app = FastAPI()

class PptxInput(BaseModel):
    file_base64: str

def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(pptx_bytes))
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    all_text.append(text)
    return "\n".join(all_text)

@app.post("/estrai-testo-pptx")
async def estrai_testo_pptx(input_data: PptxInput):
    try:
        pptx_bytes = base64.b64decode(input_data.file_base64)
        testo_slide = extract_text_from_pptx(pptx_bytes)
        return JSONResponse(content={ "testo_slide": testo_slide })
    except Exception as e:
        return JSONResponse(status_code=500, content={ "errore": str(e) })
